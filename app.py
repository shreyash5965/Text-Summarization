import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.summarize import load_summarize_chain
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
from transformers import T5Tokenizer, T5ForConditionalGeneration
from transformers import pipeline
from pptx import Presentation
import torch
import base64
import PyPDF2
from PyPDF2 import PdfWriter
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document

def create_pdf_from_text(text, output_file):
    presentation_path = output_file
    c = canvas.Canvas(presentation_path, pagesize=letter)
    
    c.setFont("Helvetica", 12)
    x = 10
    y = 700
    
    lines = text.split("\n")
    
    y = 700
    page_height = 792
    
    for line in lines:
        text_width = c.stringWidth(line)
        
        if text_width > 492:
            chunks = [line[i:i+70] for i in range(0, len(line), 70)]
            for chunk in chunks:
                c.drawString(50, y, chunk)
                y -= 20
                if y <= 50:
                    c.showPage()
                    y = page_height - 50
        else:
            c.drawString(50, y, line)
            y -= 20
        
        if y <= 50:
            c.showPage()
            y = page_height - 50

    c.save()
    return presentation_path

def create_pdf_from_text_old(text, output_file):
    output_file = output_file.replace('data/','')
    pdf_writer = PdfWriter()

    segments = [text[i:i+1000] for i in range(0, len(text), 1000)]
    i=1

    for segment in segments:
        pdf_page = pdf_writer.add_blank_page(width=800, height=792)
        pdf_page._data_stream = segment
        pdf_writer.add_page(pdf_page)

    # pdf_page = pdf_writer.add_blank_page(width=100, height=792)
    # pdf_page._data_stream = text
    # pdf_writer.add_page(pdf_page)

    with open(output_file, "wb") as pdf_output_file:
        pdf_writer.write(pdf_output_file)

def extract_text_from_doc(docx_file):
    file_path = docx_file
    doc = Document(file_path)
    extracted_text = ""
    for paragraph in doc.paragraphs:
        extracted_text += paragraph.text + "\n"
    
    return extracted_text 

#Extact the data from PPT
def extract_text_from_presentation(presentation_path):
    prs = Presentation(presentation_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

#MODEL AND TOKENIZER
checkpoint = "LaMini-Flan-T5-248M"
tokenizer = T5Tokenizer.from_pretrained(checkpoint)
base_model = T5ForConditionalGeneration.from_pretrained(checkpoint, device_map = 'auto', torch_dtype = torch.float32)

#File Loader and perprocessing
def file_preprocessing(file):
    loader = PyPDFLoader(file)
    pages = loader.load_and_split()
    text_spliiter = RecursiveCharacterTextSplitter(chunk_size = 200, chunk_overlap = 50)
    texts = text_spliiter.split_documents(pages)
    final_texts = ""

    for text in texts:
        final_texts += text.page_content

    return final_texts

#LM Pipeline
def llm_pipeline(filepath):
    pipe_sum = pipeline(
        'summarization',
        model = base_model,
        tokenizer = tokenizer,
        max_length = 500,
        min_length = 50
    )
    input_text = file_preprocessing(filepath)
    result = pipe_sum(input_text)
    result = result[0]['summary_text']
    return result

@st.cache_data
#To Display PDF of a given File

def displayPDF(file):
    #Opening file fromfile path
    with open(file, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')

    #Embedding PDF in HTNL
    pdf_display = F'<iframe src="data:application/pdf; base64, {base64_pdf}" width="100%" height="600" type="application/pdf"></iframe>'

    #Displaying File
    st.markdown(pdf_display, unsafe_allow_html=True)

# Streamlit code
st.set_page_config(layout='wide')

def openPDF(upload_file, filepath, col1, col2):
        # with open(filepath, 'wb') as temp_file:
        #     temp_file.write(upload_file.read())
        with col1:
            st.info("Uploaded File")
            pdf_viwer = displayPDF(filepath)
        with col2:
            st.info("Summarization:")
            
            summary = llm_pipeline(filepath)
            st.success(summary)

def main():

    st.title('Document Summarization App using Language Model')

    upload_file = st.file_uploader("Upload your file", type=["pdf", "docx", "pptx"])

    if upload_file is not None:
        if st.button("Summarize"):
            col1, col2 = st.columns(2)
            filepath = "data/"+upload_file.name
            print(upload_file.type)
            text_data = ""
            if upload_file.type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                text_data = extract_text_from_presentation(filepath)
                filepath = filepath.replace('pptx','pdf')
                filepath = create_pdf_from_text(text_data, filepath)
            elif upload_file.type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                text_data = extract_text_from_doc(filepath)
                filepath = filepath.replace('docx','pdf')
                filepath = create_pdf_from_text(text_data, filepath)

            openPDF(upload_file, filepath, col1, col2)

if __name__ == '__main__':
    main()